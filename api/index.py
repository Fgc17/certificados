
import logging
from flask import Flask, request, send_file
from pptx import Presentation
import pandas as pd
import io
import zipfile
from flask_cors import CORS

app = Flask(__name__)

# Set the logging level to DEBUG
app.logger.setLevel(logging.DEBUG)

# Add a StreamHandler to the logger, which will display logs in the console
stream_handler = logging.StreamHandler()
app.logger.addHandler(stream_handler)

CORS(app) 

@app.route("/api/process_files", methods=['POST'])
def process_files():

    name_index = request.form['name_index']

    cpf_index = request.form['cpf_index']

    # Get the uploaded files
    excel_file = request.files['excel']
    ppt_file = request.files['powerpoint']

    # Load the Excel file using pandas
    df = pd.read_excel(excel_file)


    # Create an in-memory zip file to store the generated PowerPoint files
    zip_output = io.BytesIO()
    zip_file = zipfile.ZipFile(zip_output, 'w', zipfile.ZIP_DEFLATED)


    for index, row in df.iterrows():
        
        # Load the PowerPoint presentation using python-pptx
        presentation = Presentation(ppt_file)

        app.logger.info(row[3])

        cpf = row[int(cpf_index) - 1]
        name = row[int(name_index) - 1]

        for slide in presentation.slides: 
            for shape in slide.shapes: 
                if not shape.has_text_frame: 
                    continue 
                for paragraph in shape.text_frame.paragraphs:  
                    for run in paragraph.runs:
                        if '###NOME###' in run.text:
                            text = run.text.replace('###NOME###', str(name))
                            run.text = ''
                            run.text = text
                        if '###CPF###' in run.text:
                            text = run.text.replace('###CPF###', str(cpf))
                            run.text = ''
                            run.text = text
        
        # Save the modified PowerPoint presentation for the current row
        slide_output = io.BytesIO()
        presentation.save(slide_output)
        slide_output.seek(0)

        # Add the PowerPoint file for the current row to the zip archive'
        zip_file.writestr(f"{name}.pptx", slide_output.getvalue())

    # Closing the ZipFile before sending it
    zip_file.close()
    # Ensuring we're at the start of the BytesIO stream
    zip_output.seek(0)

    # Create the response
    response = send_file(zip_output, mimetype='application/zip', as_attachment=True, download_name='powerpoint_files.zip')

    app.logger.info(response)

    return response

if __name__ == "__main__":
    app.run(debug=True)